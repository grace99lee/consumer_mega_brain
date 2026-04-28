from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from exporters.base import BaseExporter
from models.schemas import AnalysisResult, Review, SentimentLabel


class PowerPointExporter(BaseExporter):
    def export(
        self,
        result: AnalysisResult,
        reviews: list[Review],
        output_dir: Path,
    ) -> list[Path]:
        output_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = output_dir / "report.pptx"

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._add_title_slide(prs, result)
        self._add_summary_slide(prs, result)
        self._add_sentiment_slide(prs, result)
        self._add_themes_slide(prs, result)
        self._add_unmet_needs_slide(prs, result)
        for persona in result.personas:
            self._add_persona_slide(prs, persona, result)
        self._add_quotes_slide(prs, result)

        prs.save(str(pptx_path))
        return [pptx_path]

    def _add_title_slide(self, prs: Presentation, result: AnalysisResult):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Consumer Insights Report"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = result.query
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(0x4A, 0x4A, 0x8A)
        p2.alignment = PP_ALIGN.CENTER

        # Meta
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p3 = tf2.paragraphs[0]
        p3.text = (
            f"{result.total_reviews} reviews | "
            f"{', '.join(s.value for s in result.sources_used)} | "
            f"{result.generated_at:%Y-%m-%d}"
        )
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        p3.alignment = PP_ALIGN.CENTER

    def _add_summary_slide(self, prs: Presentation, result: AnalysisResult):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Executive Summary")

        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = result.executive_summary
        p.font.size = Pt(18)
        p.line_spacing = Pt(28)

    def _add_sentiment_slide(self, prs: Presentation, result: AnalysisResult):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Sentiment Overview")

        s = result.sentiment

        # Overall score
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Overall: {s.overall.value.upper()}"
        p.font.size = Pt(28)
        p.font.bold = True

        p2 = tf.add_paragraph()
        p2.text = f"Score: {s.score:+.2f} (scale: -1.0 to +1.0)"
        p2.font.size = Pt(18)

        # Distribution
        if s.distribution:
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(5), Inches(3))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p3 = tf2.paragraphs[0]
            p3.text = "Distribution"
            p3.font.size = Pt(20)
            p3.font.bold = True
            for label, count in s.distribution.items():
                p_item = tf2.add_paragraph()
                label_str = label.value if isinstance(label, SentimentLabel) else str(label)
                p_item.text = f"  {label_str}: {count}"
                p_item.font.size = Pt(16)

        # By source
        if s.by_source:
            txBox3 = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5), Inches(5))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            p4 = tf3.paragraphs[0]
            p4.text = "By Source"
            p4.font.size = Pt(20)
            p4.font.bold = True
            for source, score in s.by_source.items():
                p_item = tf3.add_paragraph()
                source_str = source.value if hasattr(source, 'value') else str(source)
                p_item.text = f"  {source_str}: {score:+.2f}"
                p_item.font.size = Pt(16)

    def _add_themes_slide(self, prs: Presentation, result: AnalysisResult):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Key Themes")

        y = 1.5
        for i, theme in enumerate(result.themes[:8]):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"{i + 1}. {theme.name} — {theme.review_count} mentions"
            p.font.size = Pt(16)
            p.font.bold = True

            p2 = tf.add_paragraph()
            p2.text = f"   {theme.description[:120]}"
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

            y += 0.7

    def _add_unmet_needs_slide(self, prs: Presentation, result: AnalysisResult):
        if not result.unmet_needs:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Unmet Consumer Needs")

        y = 1.5
        for need in result.unmet_needs[:6]:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"{need.need}"
            p.font.size = Pt(18)
            p.font.bold = True

            p2 = tf.add_paragraph()
            p2.text = f"Frequency: {need.frequency} | Opportunity Score: {need.opportunity_score:.2f}"
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

            y += 0.9

    def _add_persona_slide(self, prs, persona, result):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, f"Persona: {persona.name}")

        # Description
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = persona.description
        p.font.size = Pt(16)

        # Demographics & prevalence
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        if persona.demographics_hints:
            p2 = tf2.paragraphs[0]
            p2.text = f"Demographics: {persona.demographics_hints}"
            p2.font.size = Pt(14)
        p3 = tf2.add_paragraph()
        p3.text = f"Prevalence: {persona.estimated_prevalence}"
        p3.font.size = Pt(14)
        p3.font.bold = True

        # Motivations & pain points
        y = 4.0
        if persona.motivations:
            txBox3 = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(5), Inches(2))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            p4 = tf3.paragraphs[0]
            p4.text = "Motivations"
            p4.font.size = Pt(16)
            p4.font.bold = True
            for m in persona.motivations[:4]:
                pi = tf3.add_paragraph()
                pi.text = f"  • {m}"
                pi.font.size = Pt(14)

        if persona.pain_points:
            txBox4 = slide.shapes.add_textbox(Inches(7), Inches(y), Inches(5), Inches(2))
            tf4 = txBox4.text_frame
            tf4.word_wrap = True
            p5 = tf4.paragraphs[0]
            p5.text = "Pain Points"
            p5.font.size = Pt(16)
            p5.font.bold = True
            for pp in persona.pain_points[:4]:
                pi = tf4.add_paragraph()
                pi.text = f"  • {pp}"
                pi.font.size = Pt(14)

    def _add_quotes_slide(self, prs: Presentation, result: AnalysisResult):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Key Consumer Quotes")

        y = 1.5
        for q in result.key_quotes[:6]:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.9))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f'"{q.quote}"'
            p.font.size = Pt(15)
            p.font.italic = True

            p2 = tf.add_paragraph()
            p2.text = f"— {q.author or 'Anonymous'} ({q.source.value}) | {q.theme}"
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

            y += 0.95

    def _add_slide_title(self, slide, title: str):
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(11), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
